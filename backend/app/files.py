"""업로드 파일(CSV/XLSX/PDF/HWP) 파싱과 AI 매핑 결과 적용 — Streamlit import 없음."""

import io
from pathlib import Path

import pandas as pd

from backend.app.repository import 금액_컬럼들, 사업단계_옵션, 편집_컬럼순서


def 엑셀로_변환(df: pd.DataFrame) -> bytes:
    buffer = io.BytesIO()
    df.to_excel(buffer, index=False, sheet_name="사업현황")
    return buffer.getvalue()


def _업로드_원본_읽기(업로드_파일) -> pd.DataFrame:
    """형식(컬럼명·순서)을 가리지 않고 업로드된 엑셀/CSV를 그대로 읽는다.

    컬럼 매핑은 여기서 강제하지 않고 AI(ai_agent.업로드_매핑_추론)가 추론하도록 넘긴다.
    """
    파일명 = 업로드_파일.name.lower()
    if 파일명.endswith(".csv"):
        df = pd.read_csv(업로드_파일)
    else:
        df = pd.read_excel(업로드_파일)
    df = df.dropna(axis=1, how="all").dropna(axis=0, how="all")
    df.columns = [str(c).strip() for c in df.columns]
    return df.reset_index(drop=True)


def _pdf_텍스트_추출(업로드_파일, 최대글자수: int = 15000) -> str:
    from pypdf import PdfReader

    reader = PdfReader(업로드_파일)
    조각들 = [(페이지.extract_text() or "") for 페이지 in reader.pages]
    전체 = "\n".join(조각들).strip()
    if len(전체) > 최대글자수:
        전체 = 전체[:최대글자수] + "\n...(이하 생략)"
    return 전체


def _hwp_텍스트_추출(업로드_파일, 최대글자수: int = 15000) -> str:
    import io
    import tempfile
    from contextlib import closing

    from hwp5.hwp5txt import TextTransform
    from hwp5.xmlmodel import Hwp5File

    with tempfile.NamedTemporaryFile(suffix=".hwp", delete=False) as tmp:
        tmp.write(업로드_파일.getvalue())
        임시경로 = tmp.name
    try:
        출력 = io.BytesIO()
        transform = TextTransform().transform_hwp5_to_text
        with closing(Hwp5File(임시경로)) as hwp파일:
            transform(hwp파일, 출력)
        전체 = 출력.getvalue().decode("utf-8", errors="ignore").strip()
    finally:
        Path(임시경로).unlink(missing_ok=True)
    if len(전체) > 최대글자수:
        전체 = 전체[:최대글자수] + "\n...(이하 생략)"
    return 전체


def _LLM_매핑_적용(원본_df: pd.DataFrame, 매핑결과: dict) -> tuple[pd.DataFrame, list[str]]:
    """AI가 추론한 컬럼/값 매핑을 실제 데이터프레임에 적용한다.

    금액·날짜 등 실제 값 자체는 AI가 다시 받아쓰게 하지 않고(수치 오기 위험) 원본 셀 값을
    그대로 가져와 코드로만 정리한다 — AI는 "어느 컬럼이 무엇인지"만 판단한다.
    """
    경고: list[str] = []
    매핑 = (매핑결과 or {}).get("매핑") or {}
    사업단계_값매핑 = (매핑결과 or {}).get("사업단계_값매핑") or {}

    기본값 = {
        "구분": "", "업체명": "", "용역명": "", "사업구분": "", "담당자": "", "주관참여구분": "",
        "사업단계": "미분류", "진행률": 0, "시작일": None, "종료일": None,
        "계약금액": 0, "기수입금액": 0, "당해년도수입금액": 0,
    }

    결과 = pd.DataFrame(index=원본_df.index)
    for 필드 in [c for c in 편집_컬럼순서 if c != "id"]:
        원본컬럼 = 매핑.get(필드)
        if 원본컬럼 and 원본컬럼 in 원본_df.columns:
            결과[필드] = 원본_df[원본컬럼]
        else:
            결과[필드] = 기본값[필드]
            경고.append(f"'{필드}'에 해당하는 컬럼을 찾지 못해 기본값으로 채웠습니다.")

    for 컬럼 in 금액_컬럼들:
        정리값 = (
            결과[컬럼].astype(str)
            .str.replace(",", "", regex=False)
            .str.replace("원", "", regex=False)
            .str.strip()
        )
        결과[컬럼] = pd.to_numeric(정리값, errors="coerce").fillna(0)

    진행률_정리 = 결과["진행률"].astype(str).str.replace("%", "", regex=False).str.strip()
    결과["진행률"] = pd.to_numeric(진행률_정리, errors="coerce").fillna(0)

    if 사업단계_값매핑:
        결과["사업단계"] = 결과["사업단계"].astype(str).str.strip().replace(사업단계_값매핑)
    미표준_단계 = sorted(set(결과["사업단계"].astype(str)) - set(사업단계_옵션))
    if 미표준_단계:
        경고.append(f"'사업단계' 값 중 인식하지 못한 표현({', '.join(미표준_단계)})은 미분류로 대체했습니다.")
        결과["사업단계"] = 결과["사업단계"].where(결과["사업단계"].isin(사업단계_옵션), "미분류")

    for 컬럼 in ["시작일", "종료일"]:
        결과[컬럼] = pd.to_datetime(결과[컬럼], errors="coerce").dt.strftime("%Y-%m-%d")

    결과["id"] = pd.NA
    return 결과[편집_컬럼순서], 경고


_제안_기본값 = {
    "구분": "", "업체명": "", "용역명": "", "사업구분": "", "담당자": "", "주관참여구분": "",
    "사업단계": "미분류", "진행률": 0, "시작일": None, "종료일": None,
    "계약금액": 0, "기수입금액": 0, "당해년도수입금액": 0,
}


def _제안_추가행들(사업목록: list[dict]) -> pd.DataFrame:
    행들 = [{필드: 항목.get(필드, 기본값) for 필드, 기본값 in _제안_기본값.items()} for 항목 in 사업목록]
    df = pd.DataFrame(행들, columns=list(_제안_기본값.keys()))
    df["id"] = pd.NA
    return df[편집_컬럼순서]


# ---------- AI 채팅 "파일 생성"(create_file) 도구가 쓰는 포맷 변환 ----------
# ai_agent.py는 항상 텍스트만 만들어내고, 확장자에 따라 여기서 실제 바이너리로 바꾼다.

_MIME_타입_맵 = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "csv": "text/csv",
    "json": "application/json",
    "html": "text/html",
    "svg": "image/svg+xml",
    "md": "text/markdown",
    "py": "text/x-python",
    "js": "text/javascript",
    "ts": "text/typescript",
    "css": "text/css",
    "sql": "text/x-sql",
    "yaml": "text/yaml",
    "yml": "text/yaml",
    "txt": "text/plain",
}


def 파일_mime타입(파일명: str) -> str:
    확장자 = 파일명.rsplit(".", 1)[-1].lower() if "." in 파일명 else ""
    return _MIME_타입_맵.get(확장자, "application/octet-stream")


def _마크다운을_docx로(내용: str) -> bytes:
    from docx import Document

    문서 = Document()
    for 줄 in 내용.splitlines():
        줄 = 줄.rstrip()
        벗긴_줄 = 줄.lstrip()
        if not 줄.strip():
            문서.add_paragraph("")
        elif 줄.startswith("### "):
            문서.add_heading(줄[4:].strip(), level=3)
        elif 줄.startswith("## "):
            문서.add_heading(줄[3:].strip(), level=2)
        elif 줄.startswith("# "):
            문서.add_heading(줄[2:].strip(), level=1)
        elif 벗긴_줄.startswith(("- ", "* ")):
            문서.add_paragraph(벗긴_줄[2:].strip(), style="List Bullet")
        else:
            문서.add_paragraph(줄)
    버퍼 = io.BytesIO()
    문서.save(버퍼)
    return 버퍼.getvalue()


def _마크다운을_pptx로(내용: str) -> bytes:
    import re

    from pptx import Presentation

    프레젠테이션 = Presentation()
    레이아웃 = 프레젠테이션.slide_layouts[1]  # 제목 + 본문
    슬라이드_텍스트들 = re.split(r"(?m)^\s*---\s*$", 내용)

    for 슬라이드_텍스트 in 슬라이드_텍스트들:
        줄들 = [l.strip() for l in 슬라이드_텍스트.splitlines() if l.strip()]
        if not 줄들:
            continue
        제목 = 줄들[0].lstrip("#").strip()
        본문_줄들 = [l.lstrip("-* ").strip() for l in 줄들[1:]]

        슬라이드 = 프레젠테이션.slides.add_slide(레이아웃)
        슬라이드.shapes.title.text = 제목
        본문_프레임 = 슬라이드.placeholders[1].text_frame
        본문_프레임.clear()
        if 본문_줄들:
            본문_프레임.text = 본문_줄들[0]
            for 줄 in 본문_줄들[1:]:
                단락 = 본문_프레임.add_paragraph()
                단락.text = 줄

    버퍼 = io.BytesIO()
    프레젠테이션.save(버퍼)
    return 버퍼.getvalue()


def 파일_생성_바이트(파일명: str, 내용: str) -> bytes:
    """create_file 도구가 만든 텍스트 content를 파일명 확장자에 맞는 바이너리로 바꾼다.

    .xlsx는 content를 CSV로, .docx/.pptx는 마크다운으로 해석한다(정확한 규칙은
    ai_agent.py의 create_file 도구 설명에 있음). 그 외 확장자는 그대로 UTF-8 텍스트.
    """
    확장자 = 파일명.rsplit(".", 1)[-1].lower() if "." in 파일명 else ""
    if 확장자 == "xlsx":
        df = pd.read_csv(io.StringIO(내용))
        return 엑셀로_변환(df)
    if 확장자 == "docx":
        return _마크다운을_docx로(내용)
    if 확장자 == "pptx":
        return _마크다운을_pptx로(내용)
    return 내용.encode("utf-8")
